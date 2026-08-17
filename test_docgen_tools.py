"""Phase 2 smoke test: local doc-gen tools (create_docx / create_pptx).

Writes real .docx/.pptx into a temp dir, reopens them with python-docx /
python-pptx, and verifies the content matches what the handler received.
Runnable directly or via pytest.
"""

from __future__ import annotations

import tempfile
from pathlib import Path

from docx import Document
from pptx import Presentation

from core import AgentState, ChatMessage, ToolCall, ToolRegistry, new_state
from tools import register_docgen_tools


def _new_state() -> AgentState:
    return new_state(target="windows", model="LFM2.5-1.2B-Instruct")


def _dispatch(registry: ToolRegistry, name: str, args: dict) -> ChatMessage:
    return registry.dispatch(_new_state(), ToolCall(id="call_0001", name=name, arguments=args))


def _registry(tmp: Path) -> ToolRegistry:
    registry = ToolRegistry()
    register_docgen_tools(registry, output_dir=tmp)
    return registry


def test_create_docx_writes_valid_document() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        msg = _dispatch(
            registry,
            "create_docx",
            {
                "title": "Weekly Report",
                "sections": [
                    {"heading": "Progress", "body": "Shipped Phase 1."},
                    {"heading": "Next", "body": "Start Phase 2."},
                ],
            },
        )
        assert msg is not None and msg.role == "tool"
        path = Path(msg.content)
        assert path.exists() and path.suffix == ".docx"

        doc = Document(str(path))
        assert doc.paragraphs[0].text == "Weekly Report"
        headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
        assert headings == ["Progress", "Next"]
        body = " ".join(p.text for p in doc.paragraphs)
        assert "Shipped Phase 1." in body


def test_create_pptx_writes_valid_presentation() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        msg = _dispatch(
            registry,
            "create_pptx",
            {
                "title": "Deck",
                "slides": [
                    {"title": "Intro", "bullets": ["One", "Two"]},
                    {"title": "Outro", "bullets": ["Three"]},
                ],
            },
        )
        assert msg is not None and msg.role == "tool"
        path = Path(msg.content)
        assert path.exists() and path.suffix == ".pptx"

        prs = Presentation(str(path))
        assert len(prs.slides) == 2
        slides = list(prs.slides)
        assert slides[0].shapes.title.text == "Intro"
        assert slides[1].shapes.title.text == "Outro"


def test_collision_gets_unique_filename() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        args = {"title": "Report", "sections": [{"heading": "H", "body": "B"}]}
        first = Path(_dispatch(registry, "create_docx", args).content)
        second = Path(_dispatch(registry, "create_docx", args).content)
        assert first.exists() and second.exists()
        assert first != second


def test_missing_required_argument_rejected() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        msg = _dispatch(registry, "create_docx", {"sections": []})
        assert msg is not None and msg.role == "tool"
        assert "invalid arguments" in msg.content


def test_section_without_body_is_accepted() -> None:
    """body/bullets are optional now: a heading alone builds an empty section."""
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        msg = _dispatch(registry, "create_docx", {"title": "T", "sections": [{"heading": "H"}]})
        assert msg is not None and msg.role == "tool"
        assert not msg.content.startswith("error")
        doc = Document(msg.content)
        assert any(p.text == "H" for p in doc.paragraphs)


def test_single_object_sections_accepted() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        msg = _dispatch(
            registry, "create_docx", {"title": "T", "sections": {"heading": "H", "body": "B"}}
        )
        assert msg is not None and msg.role == "tool"
        assert not msg.content.startswith("error")


def test_bullets_as_string_accepted() -> None:
    """1.2B models often emit bullets as one string instead of a list."""
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        msg = _dispatch(
            registry,
            "create_pptx",
            {"title": "T", "slides": [{"title": "Intro", "bullets": "One\nTwo"}]},
        )
        assert msg is not None and msg.role == "tool"
        assert not msg.content.startswith("error")
        prs = Presentation(msg.content)
        assert len(prs.slides) == 1


def test_single_object_slides_accepted() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        msg = _dispatch(
            registry,
            "create_pptx",
            {"title": "T", "slides": {"title": "Intro", "bullets": ["One"]}},
        )
        assert msg is not None and msg.role == "tool"
        assert not msg.content.startswith("error")


def test_registry_definitions_loaded() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        registry = _registry(Path(tmp))
        names = {d.name for d in registry.definitions()}
        assert names == {"create_docx", "create_pptx"}


def main() -> None:
    tests = [v for k, v in sorted(globals().items()) if k.startswith("test_") and callable(v)]
    failures = 0
    for test in tests:
        try:
            test()
            print(f"PASS  {test.__name__}")
        except Exception as exc:  # noqa: BLE001
            failures += 1
            print(f"FAIL  {test.__name__}: {type(exc).__name__}: {exc}")
    if failures:
        raise SystemExit(f"{failures} test(s) failed")
    print(f"\nAll {len(tests)} doc-gen smoke tests passed.")


if __name__ == "__main__":
    main()