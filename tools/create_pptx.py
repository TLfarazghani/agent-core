"""create_pptx handler — Windows backend via python-pptx.

Schema (research doc §3.4): ``{ "title": str, "slides": [{title, bullets[]}] }``.
Writes the presentation to the output dir and returns the created path as the
tool result.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Callable

from pptx import Presentation

from ._paths import default_output_dir, slugify, unique_path

TOOL_NAME = "create_pptx"

_TITLE_AND_CONTENT_LAYOUT = 1


def _as_bullet_list(value: Any) -> list[str]:
    """Accept a string, a list of strings, or None; always return a list."""
    if value is None:
        return []
    if isinstance(value, str):
        lines = [line.strip() for line in value.splitlines() if line.strip()]
        return lines or [value]
    if isinstance(value, list):
        return [str(item) for item in value]
    return [str(value)]


def _normalize_slides(slides: Any) -> list[dict[str, Any]]:
    """Accept a single slide object or a list of them; fill missing fields."""
    if isinstance(slides, dict):
        slides = [slides]
    return [
        {
            "title": str(slide.get("title", "")),
            "bullets": _as_bullet_list(slide.get("bullets")),
        }
        for slide in slides
    ]


def make_handler(output_dir: Path | None = None) -> Callable[[dict[str, Any]], str]:
    out_dir = output_dir or default_output_dir()

    def handler(arguments: dict[str, Any]) -> str:
        title = arguments["title"]
        slides = _normalize_slides(arguments["slides"])

        prs = Presentation()
        layout = prs.slide_layouts[_TITLE_AND_CONTENT_LAYOUT]
        for slide_def in slides:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def["title"]
            body = slide.placeholders[1].text_frame
            for index, bullet in enumerate(slide_def["bullets"]):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = bullet

        out_dir.mkdir(parents=True, exist_ok=True)
        path = unique_path(out_dir, slugify(title), "pptx")
        prs.save(str(path))
        return str(path)

    return handler